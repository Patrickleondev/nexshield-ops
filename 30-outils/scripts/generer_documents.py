#!/usr/bin/env python3
"""Génère les gabarits DOCX et PPTX à la charte.

Usage : python3 30-outils/scripts/generer_documents.py [--out 90-templates/build]

Produit :
  - Modele-rapport.docx      : gabarit de rapport, styles définis, à remplir
  - Modele-restitution.pptx  : support de restitution client
  - Modele-presentation.pptx : présentation de la société (avant-vente)

Les styles sont définis dans le document : modifier le style dans Word met à jour
tout le document. On ne met pas en forme paragraphe par paragraphe.
"""
from __future__ import annotations
import argparse, os, sys

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
import charte as C

from docx import Document
from docx.enum.section import WD_SECTION
from docx.enum.text import WD_ALIGN_PARAGRAPH, WD_TAB_ALIGNMENT
from docx.enum.table import WD_TABLE_ALIGNMENT
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Cm, Pt, RGBColor

from pptx import Presentation
from pptx.dml.color import RGBColor as PRGB
from pptx.util import Cm as PCm, Pt as PPt


# ---------------------------------------------------------------- DOCX outils
def _fond(element, hexa: str) -> None:
    ombre = OxmlElement("w:shd")
    ombre.set(qn("w:val"), "clear")
    ombre.set(qn("w:fill"), hexa)
    element.append(ombre)


def _style(doc, nom, taille, gras, couleur, espace_avant=0, espace_apres=6,
           police=None, base="Normal"):
    styles = doc.styles
    st = styles[nom] if nom in [s.name for s in styles] else \
        styles.add_style(nom, 1)   # 1 = WD_STYLE_TYPE.PARAGRAPH
    st.font.name = police or C.POLICE_CORPS
    st.font.size = Pt(taille)
    st.font.bold = gras
    st.font.color.rgb = RGBColor(*C.rgb(couleur))
    st.paragraph_format.space_before = Pt(espace_avant)
    st.paragraph_format.space_after = Pt(espace_apres)
    # Repli de police pour les postes sans Inter
    rpr = st.element.get_or_add_rPr()
    rf = rpr.find(qn("w:rFonts"))
    if rf is None:
        rf = OxmlElement("w:rFonts"); rpr.append(rf)
    rf.set(qn("w:ascii"), police or C.POLICE_CORPS)
    rf.set(qn("w:hAnsi"), police or C.POLICE_CORPS)
    rf.set(qn("w:cs"), C.REPLI_SANS)
    return st


def _pied(section, texte_gauche: str) -> None:
    p = section.footer.paragraphs[0]
    p.text = ""
    p.paragraph_format.tab_stops.add_tab_stop(Cm(16), WD_TAB_ALIGNMENT.RIGHT)
    r = p.add_run(f"{texte_gauche}\t")
    r.font.size = Pt(8)
    r.font.color.rgb = RGBColor(*C.rgb(C.TEXTE_FAIBLE))
    r.font.name = C.POLICE_CORPS
    # Champ « page X sur Y »
    for instr in ("PAGE", "NUMPAGES"):
        if instr == "NUMPAGES":
            s = p.add_run(" / "); s.font.size = Pt(8)
            s.font.color.rgb = RGBColor(*C.rgb(C.TEXTE_FAIBLE))
        deb = OxmlElement("w:fldChar"); deb.set(qn("w:fldCharType"), "begin")
        txt = OxmlElement("w:instrText"); txt.set(qn("xml:space"), "preserve")
        txt.text = f" {instr} "
        fin = OxmlElement("w:fldChar"); fin.set(qn("w:fldCharType"), "end")
        run = p.add_run(); run.font.size = Pt(8)
        run.font.color.rgb = RGBColor(*C.rgb(C.TEXTE_FAIBLE))
        run._r.append(deb); run._r.append(txt); run._r.append(fin)


def _tableau(doc, entetes: list[str], lignes: list[list[str]], largeurs=None):
    t = doc.add_table(rows=1, cols=len(entetes))
    t.alignment = WD_TABLE_ALIGNMENT.CENTER
    t.style = "Table Grid"
    for i, e in enumerate(entetes):
        cel = t.rows[0].cells[i]
        cel.text = ""
        r = cel.paragraphs[0].add_run(e)
        r.font.bold = True
        r.font.size = Pt(9.5)
        r.font.color.rgb = RGBColor(*C.rgb(C.BLANC))
        r.font.name = C.POLICE_TITRE
        _fond(cel._tc.get_or_add_tcPr(), C.PRIMAIRE)
    for n, ligne in enumerate(lignes):
        cells = t.add_row().cells
        for i, v in enumerate(ligne):
            cells[i].text = ""
            r = cells[i].paragraphs[0].add_run(v)
            r.font.size = Pt(9.5)
            r.font.name = C.POLICE_CORPS
            r.font.color.rgb = RGBColor(*C.rgb(C.TEXTE))
            if n % 2 == 1:
                _fond(cells[i]._tc.get_or_add_tcPr(), C.FOND_ALTERNE)
    if largeurs:
        for i, l in enumerate(largeurs):
            for row in t.rows:
                row.cells[i].width = Cm(l)
    doc.add_paragraph()
    return t


# ---------------------------------------------------------------- DOCX rapport
def modele_rapport(chemin: str) -> None:
    doc = Document()

    for s in doc.sections:
        s.top_margin = s.bottom_margin = Cm(2.5)
        s.left_margin = s.right_margin = Cm(2.5)

    n = doc.styles["Normal"]
    n.font.name = C.POLICE_CORPS
    n.font.size = Pt(C.TAILLES["corps"])
    n.font.color.rgb = RGBColor(*C.rgb(C.TEXTE))
    n.paragraph_format.space_after = Pt(6)
    n.paragraph_format.line_spacing = 1.25

    _style(doc, "Heading 1", C.TAILLES["t1"], True, C.PRIMAIRE, 18, 8, C.POLICE_TITRE)
    _style(doc, "Heading 2", C.TAILLES["t2"], True, C.SECONDAIRE, 14, 6, C.POLICE_TITRE)
    _style(doc, "Heading 3", C.TAILLES["t3"], True, C.SECONDAIRE, 10, 4, C.POLICE_TITRE)
    _style(doc, "NS Legende", C.TAILLES["petit"], False, C.TEXTE_FAIBLE, 0, 10)
    _style(doc, "NS Code", C.TAILLES["mono"], False, C.TEXTE, 4, 8, C.POLICE_MONO)
    _style(doc, "NS Encadre", C.TAILLES["corps"], False, C.SECONDAIRE, 6, 10)

    # --- Couverture
    for _ in range(6):
        doc.add_paragraph()
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.LEFT
    r = p.add_run(C.SOCIETE)
    r.font.size, r.font.bold, r.font.name = Pt(34), True, C.POLICE_TITRE
    r.font.color.rgb = RGBColor(*C.rgb(C.PRIMAIRE))
    p = doc.add_paragraph()
    r = p.add_run(C.BASELINE)
    r.font.size, r.font.name = Pt(11), C.POLICE_CORPS
    r.font.color.rgb = RGBColor(*C.rgb(C.ACCENT))

    for _ in range(3):
        doc.add_paragraph()
    p = doc.add_paragraph()
    r = p.add_run("Rapport de test d'intrusion")
    r.font.size, r.font.bold, r.font.name = Pt(24), True, C.POLICE_TITRE
    r.font.color.rgb = RGBColor(*C.rgb(C.SECONDAIRE))
    p = doc.add_paragraph()
    r = p.add_run("<Périmètre de la mission>")
    r.font.size = Pt(13)
    r.font.color.rgb = RGBColor(*C.rgb(C.TEXTE_FAIBLE))

    for _ in range(2):
        doc.add_paragraph()
    _tableau(doc,
             ["", ""],
             [["Client", "<Raison sociale>"],
              ["Référence", "<AAAAMMJJ-CLIENT-RAPPORT-titre-v1.0>"],
              ["Version", "v1.0"],
              ["Date", "<AAAA-MM-JJ>"],
              ["Auteurs", "<Noms>"],
              ["Relecteurs", "<Noms>"],
              ["Doctrine appliquée", "pentest-audit v<X.Y>"],
              ["Classification", C.CLASSIFICATION]],
             largeurs=[5.5, 10.5])

    p = doc.add_paragraph(style="NS Legende")
    p.add_run("Ce document contient des informations sur les vulnérabilités des "
              "systèmes du client. Sa diffusion est strictement limitée aux "
              "destinataires désignés. Toute reproduction ou transmission non "
              "autorisée est interdite.")

    doc.add_page_break()

    # --- Sommaire (champ TOC, à mettre à jour dans Word : F9)
    doc.add_heading("Sommaire", level=1)
    p = doc.add_paragraph()
    run = p.add_run()
    deb = OxmlElement("w:fldChar"); deb.set(qn("w:fldCharType"), "begin")
    txt = OxmlElement("w:instrText"); txt.set(qn("xml:space"), "preserve")
    txt.text = r' TOC \o "1-3" \h \z \u '
    sep = OxmlElement("w:fldChar"); sep.set(qn("w:fldCharType"), "separate")
    ph = OxmlElement("w:t"); ph.text = "Clic droit > Mettre à jour les champs (F9)"
    fin = OxmlElement("w:fldChar"); fin.set(qn("w:fldCharType"), "end")
    for e in (deb, txt, sep, ph, fin):
        run._r.append(e)
    doc.add_page_break()

    def guide(texte: str) -> None:
        """Consigne de rédaction. À supprimer avant livraison au client."""
        p = doc.add_paragraph(style="NS Encadre")
        r = p.add_run("[Consigne] " + texte)
        r.italic = True
        r.font.color.rgb = RGBColor(*C.rgb(C.ACCENT))

    def champ(texte: str = "<…>") -> None:
        doc.add_paragraph(texte)

    # ============================ SYNTHÈSE EXÉCUTIVE ========================
    doc.add_heading("Synthèse exécutive", level=1)
    guide("Deux pages maximum. Aucun jargon. Lisible par un directeur général "
          "qui n'ouvrira pas le reste du document. Si un terme technique est "
          "indispensable, l'expliquer en une ligne.")

    doc.add_heading("Contexte et objectif", level=2)
    champ("<Pourquoi cette mission, à la demande de qui, dans quel cadre "
          "(obligation réglementaire, exigence d'un client, incident, projet).>")

    doc.add_heading("Verdict", level=2)
    guide("Une phrase, orientée impact métier. Exemple : « Un attaquant "
          "disposant d'un simple accès Internet peut prendre le contrôle "
          "complet de l'application et accéder aux données de l'ensemble des "
          "clients, en moins de deux heures et sans compte valide. »")
    champ()

    doc.add_heading("Constatations en un coup d'œil", level=2)
    _tableau(doc, ["Sévérité", "Nombre", "Délai de correction recommandé"],
             [["Critique", "0", "Immédiat (< 72 h)"],
              ["Élevée", "0", "30 jours"],
              ["Moyenne", "0", "90 jours"],
              ["Faible", "0", "Prochain cycle de maintenance"],
              ["Information", "0", "Aucune obligation"]],
             largeurs=[4, 3, 9])

    doc.add_heading("Les trois choses à corriger en premier", level=2)
    for i in range(1, 4):
        doc.add_paragraph(
            f"<Titre {i}> — <impact métier en une phrase> → <action> — "
            f"<effort estimé>", style="List Number")

    doc.add_heading("Appréciation globale", level=2)
    guide("Posture générale, et les points forts constatés — il y en a "
          "toujours, et les citer crédibilise les critiques. Tendance par "
          "rapport à un éventuel audit précédent.")
    champ()

    doc.add_page_break()

    # ============================ 1. CADRE ==================================
    doc.add_heading("1. Cadre de la mission", level=1)

    doc.add_heading("1.1 Périmètre testé", level=2)
    _tableau(doc, ["#", "Actif", "Type", "Identifiant", "Environnement"],
             [["1", "<…>", "<web / API / IP / mobile / LLM>", "<…>", "<prod / préprod>"],
              ["2", "", "", "", ""]],
             largeurs=[1.2, 4, 3.5, 4.3, 3])
    guide("Repris du RoE §2. En cas d'écart avec le RoE, le justifier ici "
          "explicitement — un écart non justifié est une faute contractuelle.")

    doc.add_heading("1.2 Périmètre exclu", level=2)
    champ("<Systèmes voisins non testés, et pourquoi.>")

    doc.add_heading("1.3 Cadre contractuel et légal", level=2)
    _tableau(doc, ["Élément", "Référence"],
             [["Accord de confidentialité", "<référence, date>"],
              ["Contrat-cadre / bon de commande", "<référence, date>"],
              ["Règles d'engagement (RoE)", "<référence, version, date>"],
              ["Autorisation de test signée", "<signataire, fonction, date>"],
              ["Autorisation de l'hébergeur", "<référence ou « non requise »>"],
              ["Droit applicable", "<Togo — loi n° 2018-026 / autre>"]],
             largeurs=[6, 10])
    guide("Cette section n'est pas administrative : elle établit que les tests "
          "étaient licites. Voir 00-societe/juridique/CADRE-LEGAL.md §1.1.")

    doc.add_heading("1.4 Méthodologie et référentiels", level=2)
    doc.add_paragraph(
        "Mission conduite selon le standard PTES, avec les tests de l'OWASP Web "
        "Security Testing Guide, sous la version <X.Y> de notre doctrine "
        "« pentest-audit ».")
    _tableau(doc, ["Référentiel", "Version", "Usage dans la mission"],
             [["PTES", "—", "Structuration des phases"],
              ["OWASP WSTG", "<v4.2>", "Catalogue de tests, preuve de couverture"],
              ["OWASP ASVS", "<v4.0.3>", "Niveau d'assurance visé : <L1 / L2 / L3>"],
              ["NIST SP 800-115", "—", "Cadre technique de référence"],
              ["MITRE ATT&CK", "<v15>", "Qualification des techniques"],
              ["CVSS", "v4.0", "Scoring des vulnérabilités"],
              ["CWE", "<v4.14>", "Classification des faiblesses"]],
             largeurs=[4.5, 3, 8.5])

    doc.add_heading("1.5 Phases réalisées", level=2)
    _tableau(doc, ["Phase PTES", "Période", "Contenu"],
             [[n, "<du … au …>", "<…>"] for n in
              ["Pré-engagement", "Renseignement", "Modélisation de menaces",
               "Analyse de vulnérabilités", "Exploitation", "Post-exploitation",
               "Restitution"]],
             largeurs=[4.5, 3.5, 8])

    doc.add_heading("1.6 Équipe et calendrier", level=2)
    _tableau(doc, ["Nom", "Rôle", "Période d'intervention"],
             [["<…>", "Chef de mission", "<…>"],
              ["<…>", "Testeur", "<…>"],
              ["<…>", "Relecteur qualité", "<…>"]],
             largeurs=[5, 5, 6])

    doc.add_heading("1.7 Limites de l'évaluation", level=2)
    guide("Section de protection — ne jamais la supprimer, même vide.")
    champ("<Ce qui n'a pas pu être testé et pourquoi : indisponibilité d'un "
          "environnement, fenêtre trop courte, fonctionnalité non déployée, "
          "compte non fourni, refus d'une technique.>")
    doc.add_paragraph(
        "Un test d'intrusion est une évaluation à un instant donné, sur un "
        "périmètre délimité. L'absence de constatation sur un composant ne "
        "garantit pas son absence de vulnérabilité.")

    doc.add_heading("1.8 Échelle de sévérité", level=2)
    _tableau(doc, ["Sévérité", "CVSS v4.0", "Définition"],
             [["Critique", "9.0 – 10.0", "Compromission immédiate, sans "
               "prérequis, avec impact majeur sur les données ou le service"],
              ["Élevée", "7.0 – 8.9", "Compromission possible avec un prérequis "
               "réaliste"],
              ["Moyenne", "4.0 – 6.9", "Impact limité ou exploitation "
               "conditionnée"],
              ["Faible", "0.1 – 3.9", "Impact marginal, à corriger par hygiène"],
              ["Information", "0.0", "Observation sans impact de sécurité direct"]],
             largeurs=[3, 3, 10])
    doc.add_paragraph(
        "La sévérité affichée tient compte de l'exposition réelle et de la "
        "valeur métier de l'actif. Toute divergence avec le score CVSS est "
        "justifiée par écrit dans la fiche de vulnérabilité concernée.")

    doc.add_page_break()

    # ============================ 2. CONSTATATIONS ==========================
    doc.add_heading("2. Constatations détaillées", level=1)
    guide("Une sous-section par vulnérabilité, par sévérité décroissante. "
          "Dupliquer le bloc ci-dessous autant que nécessaire.")

    doc.add_heading("2.1 <CLIENT>-2026-001 — <Titre orienté impact, pas outil>", level=2)
    _tableau(doc, ["Champ", "Valeur"],
             [["Identifiant", "<CLIENT>-2026-001"],
              ["Sévérité", "Critique"],
              ["Score CVSS v4.0", "9.3"],
              ["Vecteur CVSS", "CVSS:4.0/AV:N/AC:L/AT:N/PR:N/UI:N/VC:H/VI:H/VA:H"],
              ["Criticité métier", "<si différente du CVSS : justifier ici>"],
              ["Actif affecté", "<…>"],
              ["CWE", "CWE-<000> — <intitulé>"],
              ["OWASP WSTG", "WSTG-<CAT>-<NN>"],
              ["OWASP ASVS", "V<x.y.z>"],
              ["MITRE ATT&CK", "T<NNNN> — <intitulé>"],
              ["Statut", "Ouverte"],
              ["Découverte par", "<nom>"],
              ["Date de découverte", "<AAAA-MM-JJ>"],
              ["Notifiée au client le", "<AAAA-MM-JJ, si critique>"]],
             largeurs=[5.5, 10.5])

    doc.add_heading("Description", level=3)
    champ("<Ce qu'est le problème, en clair, avant tout détail technique.>")

    doc.add_heading("Conditions d'exploitation", level=3)
    _tableau(doc, ["Prérequis", "Valeur"],
             [["Position de l'attaquant", "<Internet / réseau interne / compte utilisateur>"],
              ["Authentification requise", "<aucune / utilisateur / administrateur>"],
              ["Interaction utilisateur", "<aucune / requise>"],
              ["Complexité", "<faible / élevée>"],
              ["Temps nécessaire", "<…>"]],
             largeurs=[6, 10])

    doc.add_heading("Impact métier", level=3)
    guide("La section qui justifie la facture. Pas « injection SQL possible » "
          "mais « accès en lecture et écriture à l'intégralité de la base "
          "clients, soit <N> personnes, incluant <types de données> ».")
    champ()

    doc.add_heading("Preuve de concept", level=3)
    guide("Reproductible pas à pas par l'équipe du client. Aucune donnée "
          "personnelle réelle : anonymiser systématiquement. Les preuves brutes "
          "restent au coffre chiffré, seules leurs empreintes figurent en "
          "annexe C.")
    doc.add_paragraph("Étape 1 — <action>", style="List Number")
    p = doc.add_paragraph(style="NS Code")
    p.add_run("<requête, commande ou charge utile>")
    doc.add_paragraph("Étape 2 — <action>", style="List Number")
    p = doc.add_paragraph(style="NS Code")
    p.add_run("<réponse observée, tronquée et anonymisée>")
    doc.add_paragraph("Capture 1 — <légende décrivant ce que la capture prouve>",
                      style="NS Legende")

    doc.add_heading("Recommandation", level=3)
    _tableau(doc, ["Horizon", "Action", "Effort", "Référence"],
             [["Mesure d'urgence", "<contournement : règle WAF, désactivation>",
               "<jours>", "<…>"],
              ["Correctif de fond", "<correction durable>", "<semaines>",
               "<OWASP Cheat Sheet, ASVS V…>"],
              ["Prévention", "<test automatisé, revue de code, formation>",
               "<…>", "<…>"]],
             largeurs=[3.5, 6, 2.5, 4])
    doc.add_paragraph(
        "Une mesure d'urgence n'est jamais présentée comme une correction : "
        "elle réduit l'exposition en attendant le correctif de fond.")

    doc.add_heading("Détection", level=3)
    guide("Notre différenciateur : comment le client détecte cette attaque "
          "s'il la subit demain. Journaux à surveiller, indicateurs, règle "
          "SIGMA proposée.")
    champ()

    doc.add_heading("Références externes", level=3)
    champ("<CVE, avis éditeur, publication de recherche.>")

    doc.add_page_break()

    # ============================ 3 à 5 =====================================
    doc.add_heading("3. Constatations hors référentiel", level=1)
    guide("Ce que la checklist ne prévoyait pas : logique métier, "
          "enchaînements de failles mineures en compromission majeure, "
          "observations d'architecture. Une section vide sur plusieurs missions "
          "d'affilée est un signal d'alerte interne, pas un bon résultat.")
    champ()

    doc.add_heading("4. Chemin d'attaque", level=1)
    guide("Le récit de la compromission, du point d'entrée à l'objectif. C'est "
          "ce que la direction retiendra. Insérer un schéma : il vaut mieux "
          "qu'une page de texte.")
    p = doc.add_paragraph(style="NS Code")
    p.add_run("Internet → <2.1> → accès applicatif → <2.3> → compte "
              "administrateur → base de données")
    doc.add_paragraph("Figure 1 — Chemin de compromission", style="NS Legende")

    doc.add_heading("5. Plan de remédiation", level=1)
    _tableau(doc, ["#", "Constatation", "Sévérité", "Action", "Effort",
                   "Délai", "Responsable"],
             [["1", "<CLIENT>-2026-001", "Critique", "<…>", "<…>", "< 72 h", "<…>"],
              ["2", "", "", "", "", "", ""]],
             largeurs=[1, 3.4, 2.2, 3.4, 1.8, 2, 2.2])
    doc.add_paragraph(
        "Contre-vérification : <prévue le … / non prévue au contrat>. Elle "
        "porte uniquement sur les constatations du présent rapport et ne "
        "constitue pas une nouvelle évaluation.")

    doc.add_page_break()

    # ============================ ANNEXES ===================================
    doc.add_heading("Annexe A — Couverture des tests", level=1)
    guide("La preuve de sérieux. La troisième colonne est la plus importante : "
          "un test non exécuté est déclaré, avec son motif.")
    _tableau(doc, ["Identifiant", "Intitulé", "Statut", "Motif si non exécuté"],
             [["WSTG-INFO-01", "<…>", "Exécuté", ""],
              ["WSTG-ATHN-04", "<…>", "Non applicable", "Pas d'authentification"],
              ["WSTG-BUSL-03", "<…>", "Non exécuté", "Fenêtre insuffisante"]],
             largeurs=[3.5, 6, 3, 3.5])
    doc.add_paragraph("Bilan : <N> exécutés, <N> non applicables, <N> non "
                      "exécutés, sur <N> tests du référentiel.")

    doc.add_heading("Annexe B — Conformité et rattachement réglementaire", level=1)
    guide("Section à conserver pour les clients soumis à une obligation. "
          "Elle transforme le rapport en pièce d'audit.")
    _tableau(doc, ["Constatation", "ISO 27001 (Annexe A)", "Exigence locale",
                   "RGPD", "Commentaire"],
             [["<CLIENT>-2026-001", "A.8.<nn>", "<Loi 2018-026 / règles ANCy>",
               "<Art. 32>", "<…>"],
              ["", "", "", "", ""]],
             largeurs=[3.2, 3.2, 3.6, 2.4, 3.6])
    doc.add_paragraph(
        "Ce tableau est indicatif. Il identifie les mesures de sécurité "
        "concernées ; il ne constitue ni un audit de conformité, ni un avis "
        "juridique.")

    doc.add_heading("Annexe C — Outillage", level=1)
    _tableau(doc, ["Outil", "Version", "Usage", "Phase"],
             [["<…>", "<…>", "<…>", "<…>"]],
             largeurs=[4, 3, 6, 3])

    doc.add_heading("Annexe D — Manifeste des preuves", level=1)
    doc.add_paragraph(
        "Les preuves sont conservées chiffrées, hors du présent document. "
        "Destruction prévue le <date de livraison + 90 jours>, actée par un "
        "certificat de destruction remis au client.")
    _tableau(doc, ["Fichier", "Description", "SHA-256"],
             [["<…>", "<…>", "<…>"]],
             largeurs=[4.5, 5.5, 6])

    doc.add_heading("Annexe E — Journal des opérations", level=1)
    guide("Extrait du journal de mission. En cas de litige sur une "
          "indisponibilité, c'est cette annexe qui établit ce que nous avons "
          "fait, et quand.")
    _tableau(doc, ["Date et heure", "Opérateur", "IP source", "Actif", "Action"],
             [["<…>", "<…>", "<…>", "<…>", "<…>"]],
             largeurs=[3, 2.8, 2.8, 3.4, 4])

    doc.add_heading("Annexe F — Glossaire", level=1)
    _tableau(doc, ["Terme", "Définition"],
             [["CVSS", "Système standardisé de notation de la gravité "
               "technique d'une vulnérabilité, de 0 à 10."],
              ["CWE", "Classification des types de faiblesses logicielles."],
              ["MITRE ATT&CK", "Base de connaissances des techniques employées "
               "par les attaquants réels."],
              ["<…>", "<…>"]],
             largeurs=[4, 12])

    doc.add_heading("Annexe G — Diffusion et confidentialité", level=1)
    _tableau(doc, ["Destinataire", "Fonction", "Date de remise", "Format"],
             [["<…>", "<…>", "<…>", "<PDF chiffré>"]],
             largeurs=[4.5, 4.5, 3.5, 3.5])
    doc.add_paragraph(
        "Ce document est classé " + C.CLASSIFICATION + ". Toute diffusion hors "
        "de la liste ci-dessus requiert l'accord écrit du client. Il contient "
        "des informations dont la divulgation faciliterait une attaque.")

    _pied(doc.sections[0], f"<CLIENT> — {C.CLASSIFICATION}")
    doc.save(chemin)
    print("écrit", chemin)


# ---------------------------------------------------------------- PPTX outils
def _diapo_vide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _zone(diapo, x, y, l, h, texte, taille, gras=False, couleur=None,
          police=None, interligne=1.15):
    tb = diapo.shapes.add_textbox(PCm(x), PCm(y), PCm(l), PCm(h))
    tf = tb.text_frame
    tf.word_wrap = True
    lignes = texte.split("\n")
    for i, ligne in enumerate(lignes):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = interligne
        r = p.add_run(); r.text = ligne
        r.font.size = PPt(taille)
        r.font.bold = gras
        r.font.name = police or C.POLICE_TITRE
        r.font.color.rgb = PRGB(*C.rgb(couleur or C.TEXTE))
    return tb


def _bandeau(diapo, prs, hauteur=1.6, couleur=None):
    from pptx.enum.shapes import MSO_SHAPE
    f = diapo.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                               prs.slide_width, PCm(hauteur))
    f.fill.solid(); f.fill.fore_color.rgb = PRGB(*C.rgb(couleur or C.PRIMAIRE))
    f.line.fill.background()
    f.shadow.inherit = False
    return f


def _filet_accent(diapo, x, y, largeur=3.0):
    from pptx.enum.shapes import MSO_SHAPE
    f = diapo.shapes.add_shape(MSO_SHAPE.RECTANGLE, PCm(x), PCm(y),
                               PCm(largeur), PCm(0.14))
    f.fill.solid(); f.fill.fore_color.rgb = PRGB(*C.rgb(C.ACCENT))
    f.line.fill.background(); f.shadow.inherit = False


def _titre_diapo(diapo, prs, titre, sur_titre=None):
    _bandeau(diapo, prs)
    if sur_titre:
        _zone(diapo, 1.6, 0.28, 20, 0.6, sur_titre.upper(), 9, False, C.ACCENT)
    _zone(diapo, 1.6, 0.68, 22, 1.0, titre, 20, True, C.BLANC)


def _notes(diapo, texte: str) -> None:
    """Note de présentateur : ce qu'on dit, pas ce qu'on affiche."""
    diapo.notes_slide.notes_text_frame.text = texte


def _puces(diapo, x, y, l, items, taille=14, interligne=1.5):
    tb = diapo.shapes.add_textbox(PCm(x), PCm(y), PCm(l), PCm(10))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = interligne
        r = p.add_run(); r.text = "—   " + item
        r.font.size = PPt(taille)
        r.font.name = C.POLICE_CORPS
        r.font.color.rgb = PRGB(*C.rgb(C.TEXTE))
    return tb


def _carte(diapo, x, y, l, h, titre, corps, couleur=None):
    """Encadré titre + texte, à la charte."""
    from pptx.enum.shapes import MSO_SHAPE
    f = diapo.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               PCm(x), PCm(y), PCm(l), PCm(h))
    f.fill.solid(); f.fill.fore_color.rgb = PRGB(*C.rgb(C.FOND_ALTERNE))
    f.line.color.rgb = PRGB(*C.rgb(couleur or C.ACCENT))
    f.line.width = PPt(1.5)
    f.shadow.inherit = False
    f.text_frame.text = ""
    _zone(diapo, x + 0.5, y + 0.4, l - 1, 0.8, titre, 12, True,
          couleur or C.PRIMAIRE)
    _zone(diapo, x + 0.5, y + 1.3, l - 1, h - 1.7, corps, 10.5, False,
          C.TEXTE_FAIBLE, C.POLICE_CORPS, 1.25)
    return f


def _histogramme_severites(diapo, x, y, largeur=22.0, hauteur=7.0):
    """Barres horizontales de répartition par sévérité, tracées en formes."""
    from pptx.enum.shapes import MSO_SHAPE
    pas = hauteur / len(C.ORDRE_SEVERITES)
    for i, sev in enumerate(C.ORDRE_SEVERITES[:5]):
        yy = y + i * pas
        _zone(diapo, x, yy, 4.2, pas, sev, 11, True,
              C.SEVERITES[sev]["fond"], C.POLICE_CORPS)
        # Barre témoin : la longueur se règle à la main lors du remplissage.
        b = diapo.shapes.add_shape(MSO_SHAPE.RECTANGLE, PCm(x + 4.4),
                                   PCm(yy + 0.12), PCm(2.0), PCm(pas - 0.5))
        b.fill.solid(); b.fill.fore_color.rgb = PRGB(*C.rgb(C.SEVERITES[sev]["fond"]))
        b.line.fill.background(); b.shadow.inherit = False
        _zone(diapo, x + 6.7, yy, 2, pas, "<n>", 12, True, C.TEXTE,
              C.POLICE_MONO)


def _chemin_attaque(diapo, x, y, etapes: list[str]):
    """Enchaînement de blocs reliés par des flèches."""
    from pptx.enum.shapes import MSO_SHAPE
    largeur, ecart = 4.6, 1.1
    for i, etape in enumerate(etapes):
        xx = x + i * (largeur + ecart)
        couleur = C.SEVERITES["Critique"]["fond"] if i == len(etapes) - 1 else C.SECONDAIRE
        b = diapo.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   PCm(xx), PCm(y), PCm(largeur), PCm(2.4))
        b.fill.solid(); b.fill.fore_color.rgb = PRGB(*C.rgb(couleur))
        b.line.fill.background(); b.shadow.inherit = False
        tf = b.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run(); r.text = etape
        r.font.size = PPt(10.5); r.font.bold = True
        r.font.name = C.POLICE_CORPS
        r.font.color.rgb = PRGB(*C.rgb(C.BLANC))
        if i < len(etapes) - 1:
            fl = diapo.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                        PCm(xx + largeur + 0.15),
                                        PCm(y + 0.9), PCm(ecart - 0.3), PCm(0.6))
            fl.fill.solid(); fl.fill.fore_color.rgb = PRGB(*C.rgb(C.ACCENT))
            fl.line.fill.background(); fl.shadow.inherit = False


def _couverture(prs, titre, sous_titre, mentions):
    d = _diapo_vide(prs)
    from pptx.enum.shapes import MSO_SHAPE
    f = d.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    f.fill.solid(); f.fill.fore_color.rgb = PRGB(*C.rgb(C.PRIMAIRE))
    f.line.fill.background(); f.shadow.inherit = False
    _zone(d, 2.2, 4.2, 20, 1.6, C.SOCIETE, 40, True, C.BLANC)
    _zone(d, 2.2, 6.1, 20, 0.8, C.BASELINE, 12, False, C.ACCENT)
    _filet_accent(d, 2.2, 7.3, 4.0)
    _zone(d, 2.2, 8.0, 20, 1.4, titre, 22, True, C.BLANC)
    _zone(d, 2.2, 9.5, 20, 0.8, sous_titre, 13, False, "94A3B8")
    _zone(d, 2.2, 12.6, 20, 1.6, mentions, 9, False, "64748B")
    return d


# ---------------------------------------------------------------- PPTX
def modele_restitution(chemin: str) -> None:
    prs = Presentation()
    prs.slide_width, prs.slide_height = PCm(33.87), PCm(19.05)   # 16:9

    _couverture(prs, "Restitution de mission",
                "<Client> — <Type de mission> — <Date>",
                f"{C.CLASSIFICATION}\nDiffusion limitée aux destinataires désignés")

    _notes(prs.slides[0],
           "Ne pas commencer par les failles. Commencer par remercier l'équipe "
           "technique du client pour sa disponibilité — c'est elle qui vous a "
           "ouvert les accès, et c'est elle qui portera les corrections.\n\n"
           "Annoncer la durée (45 min de présentation, 15 min de questions) et "
           "préciser que le rapport détaillé a déjà été transmis.")

    # 1. Déroulé
    d = _diapo_vide(prs)
    _titre_diapo(d, prs, "Ce que nous allons voir", "Déroulé")
    _puces(d, 1.8, 3.4, 20, [
        "Le cadre : ce qui a été testé, et ce qui ne l'a pas été",
        "Le verdict, en une phrase",
        "Le chemin d'attaque : comment nous sommes entrés",
        "Les constatations, par ordre de priorité",
        "Le plan de remédiation",
        "Vos questions",
    ], 15)
    _notes(d, "20 secondes sur cette diapo. Elle sert seulement à rassurer : "
              "le client sait combien de temps ça va durer et quand il pourra "
              "poser ses questions.")

    # 2. Cadre
    d = _diapo_vide(prs)
    _titre_diapo(d, prs, "Ce qui a été testé", "Cadre de la mission")
    _carte(d, 1.8, 3.2, 9.6, 5.4, "Dans le périmètre",
           "<Actifs testés>\n<Environnement>\n<Période>", C.ACCENT)
    _carte(d, 12.2, 3.2, 9.6, 5.4, "Hors périmètre",
           "<Ce qui n'a pas été testé, et pourquoi>", C.TEXTE_FAIBLE)
    _carte(d, 22.6, 3.2, 9.6, 5.4, "Limites",
           "<Fenêtre, accès manquants, fonctionnalités non déployées>",
           C.SEVERITES["Moyenne"]["fond"])
    _zone(d, 1.8, 9.4, 30, 2,
          "Méthodologie : PTES, OWASP WSTG <version>, doctrine interne v<X.Y>.\n"
          "Chaque constatation est rattachée à un identifiant MITRE ATT&CK.",
          12, False, C.TEXTE_FAIBLE, C.POLICE_CORPS)
    _notes(d, "Insister sur la colonne « Limites ». Un client qui découvre après "
              "coup qu'une zone n'a pas été testée perd confiance dans tout le "
              "reste. Le dire soi-même, en premier, protège la crédibilité du "
              "rapport.\n\nSi le périmètre a été réduit en cours de mission, "
              "c'est ici qu'on l'explique — pas dans les questions.")

    # 3. Verdict
    d = _diapo_vide(prs)
    _titre_diapo(d, prs, "Le verdict", "Synthèse")
    _zone(d, 2.4, 6.0, 29, 5,
          "« <Un attaquant disposant de … peut … en moins de … > »",
          26, True, C.PRIMAIRE, C.POLICE_TITRE, 1.3)
    _filet_accent(d, 2.4, 11.4, 6.0)
    _notes(d, "Une seule phrase, apprise par cœur. C'est LA phrase que le "
              "dirigeant répétera à son conseil d'administration.\n\n"
              "Formuler en impact métier, jamais en vocabulaire technique. "
              "Marquer un silence de deux secondes après l'avoir dite.")

    # 4. Répartition
    d = _diapo_vide(prs)
    _titre_diapo(d, prs, "Répartition des constatations", "Résultats")
    _histogramme_severites(d, 2.0, 3.6, hauteur=8.0)
    _zone(d, 14.0, 3.6, 18, 8,
          "Ajuster la longueur des barres à la main après remplissage.\n\n"
          "Rappel de lecture :\n"
          "La sévérité tient compte de l'exposition réelle et de la valeur\n"
          "métier de l'actif, pas seulement du score CVSS.",
          12, False, C.TEXTE_FAIBLE, C.POLICE_CORPS, 1.4)
    _notes(d, "Ne pas s'attarder sur les chiffres. Le nombre de failles n'a "
              "aucune importance : c'est la diapo suivante qui compte.\n\n"
              "Si le client demande « c'est beaucoup ou pas ? », répondre "
              "honnêtement par comparaison avec des systèmes équivalents, sans "
              "dramatiser ni minimiser.")

    # 5. Chemin d'attaque
    d = _diapo_vide(prs)
    _titre_diapo(d, prs, "Comment nous sommes entrés", "Chemin d'attaque")
    _chemin_attaque(d, 1.8, 6.4,
                    ["Internet\nsans compte", "<Faille 1>\naccès applicatif",
                     "<Faille 2>\nélévation", "Administrateur",
                     "Base de données\nclients"])
    _zone(d, 1.8, 10.4, 30, 2,
          "Durée réelle de la chaîne complète : <…>. Aucune alerte n'a été "
          "déclenchée côté client pendant cette progression.",
          12, False, C.TEXTE_FAIBLE, C.POLICE_CORPS)
    _notes(d, "La diapo la plus importante de la présentation. Raconter, ne pas "
              "lister. Prendre 3 à 4 minutes.\n\n"
              "La phrase du bas — « aucune alerte n'a été déclenchée » — est "
              "celle qui ouvre la discussion sur la détection, donc sur la "
              "mission suivante. Ne pas la dire de façon accusatrice : c'est un "
              "constat, pas un reproche.")

    # 6-7. Constatations
    for n in (1, 2):
        d = _diapo_vide(prs)
        _titre_diapo(d, prs, f"<CLIENT>-2026-00{n} — <titre>",
                     f"Constatation {n}")
        _carte(d, 1.8, 3.2, 14.6, 4.4, "Impact métier",
               "<Ce que ça permet concrètement, chiffré>",
               C.SEVERITES["Critique"]["fond"])
        _carte(d, 17.4, 3.2, 14.6, 4.4, "Mesure d'urgence",
               "<Contournement applicable en quelques jours>",
               C.SEVERITES["Moyenne"]["fond"])
        _carte(d, 1.8, 8.2, 14.6, 4.4, "Correctif de fond",
               "<Correction durable, effort estimé>", C.ACCENT)
        _carte(d, 17.4, 8.2, 14.6, 4.4, "Détection",
               "<Journaux à surveiller, règle proposée>",
               C.SEVERITES["Corrigée"]["fond"])
        _zone(d, 1.8, 13.2, 30, 1.2,
              "CVSS v4.0 <score>  ·  CWE-<nnn>  ·  WSTG-<xxx>  ·  ATT&CK T<nnnn>",
              11, False, C.TEXTE_FAIBLE, C.POLICE_MONO)
        _notes(d, "Une constatation par diapo, trois au maximum dans toute la "
                  "présentation. Le reste est dans le rapport.\n\n"
                  "Ne jamais lire la carte « Impact métier » : la dire avec ses "
                  "mots, en regardant le dirigeant. Les trois autres cartes "
                  "s'adressent à l'équipe technique.")

    # 8. Remédiation
    d = _diapo_vide(prs)
    _titre_diapo(d, prs, "Par où commencer", "Plan de remédiation")
    _carte(d, 1.8, 3.2, 9.6, 8.0, "Sous 72 heures",
           "<Actions critiques>\n\nEffort : <…>",
           C.SEVERITES["Critique"]["fond"])
    _carte(d, 12.2, 3.2, 9.6, 8.0, "Sous 30 jours",
           "<Actions de sévérité élevée>\n\nEffort : <…>",
           C.SEVERITES["Élevée"]["fond"])
    _carte(d, 22.6, 3.2, 9.6, 8.0, "Sous 90 jours",
           "<Actions de sévérité moyenne>\n\nEffort : <…>",
           C.SEVERITES["Moyenne"]["fond"])
    _notes(d, "Vérifier avant la réunion que les efforts annoncés sont "
              "réalistes au regard des moyens du client. Une recommandation "
              "impossible à tenir sera ignorée, et discréditera les autres.\n\n"
              "Demander explicitement : « Est-ce que ce calendrier vous paraît "
              "tenable ? » La réponse vaut plus que le plan.")

    # 9. Points forts
    d = _diapo_vide(prs)
    _titre_diapo(d, prs, "Ce qui fonctionne déjà", "Points forts")
    _puces(d, 1.8, 3.6, 28, ["<Point fort 1>", "<Point fort 2>", "<Point fort 3>"], 15)
    _notes(d, "Diapo obligatoire, jamais supprimée. Il y a toujours des points "
              "forts.\n\nDeux raisons : elle crédibilise les critiques (on n'est "
              "pas là pour tout démolir) et elle donne à l'équipe technique du "
              "client de quoi défendre son travail devant sa direction. C'est ce "
              "qui en fait un allié pour la suite.")

    # 10. Suite
    d = _diapo_vide(prs)
    _titre_diapo(d, prs, "Ce que nous proposons", "Suite")
    _puces(d, 1.8, 3.6, 28, [
        "Contre-vérification à J+30 sur les constatations corrigées",
        "Accompagnement de vos équipes à la correction",
        "Mise en détection des scénarios identifiés",
        "Destruction des preuves à J+90, avec certificat",
    ], 15)
    _notes(d, "Ne pas vendre agressivement. Présenter ces options comme la "
              "suite logique, et laisser le client demander un devis.\n\n"
              "Insister en revanche sur la destruction des preuves à J+90 : "
              "c'est un engagement, pas une option, et c'est rassurant.")

    # 11. Questions
    d = _diapo_vide(prs)
    _titre_diapo(d, prs, "Vos questions", "")
    _zone(d, 1.8, 8.0, 28, 3,
          "<Contact du chef de mission — téléphone et courriel>",
          14, False, C.TEXTE_FAIBLE, C.POLICE_CORPS)
    _notes(d, "Questions difficiles à préparer :\n\n"
              "« Est-ce qu'on est en règle ? » → Ce n'est pas un audit de "
              "conformité. On peut dire quelles mesures sont concernées, pas "
              "délivrer un avis juridique.\n\n"
              "« Est-ce qu'on a déjà été piraté ? » → Notre mission ne portait "
              "pas là-dessus. Si le sujet inquiète, c'est une prestation "
              "distincte (recherche de compromission).\n\n"
              "« Vous garantissez qu'il n'y a plus rien ? » → Non, et personne "
              "ne peut le garantir. Évaluation à un instant donné, sur un "
              "périmètre délimité.")

    prs.save(chemin)
    print("écrit", chemin, f"({len(prs.slides._sldIdLst)} diapos)")


def modele_presentation(chemin: str) -> None:
    prs = Presentation()
    prs.slide_width, prs.slide_height = PCm(33.87), PCm(19.05)

    _couverture(prs, "Présentation de la société",
                "Sécurité offensive, IA et conformité",
                "<Date> — <Interlocuteur>")

    plan = [
        ("Qui nous sommes", "Une équipe, pas un intermédiaire",
         "<N> associés, <N> ans d'expérience cumulée\n"
         "Certifications : <…>\n"
         "Nos travaux publics : <writeups, CVE, outils>"),
        ("Le constat", "Pourquoi nous existons",
         "Trop d'organisations sont traitées comme des cibles commerciales\n"
         "plutôt que comme des partenaires de sécurité.\n\n"
         "Notre approche : honnête, offensive d'abord, fondée sur les menaces\n"
         "réelles — pas sur des listes de conformité."),
        ("Nos services", "Huit offres, une seule doctrine",
         "Pentest & audit\nAI RedTeaming\nSécurité applicative\nDevSecOps\n"
         "SOC & outillage IA défensif\nX-Privacy\nSensibilisation\n"
         "Infrastructure, VPN & Cloudflare"),
        ("Notre méthode", "Des référentiels reconnus, appliqués de bout en bout",
         "ISO/IEC 27001 — notre système de management\n"
         "MITRE ATT&CK — le langage commun de tous nos livrables\n"
         "PTES, OWASP WSTG, ASVS — l'exécution\n"
         "OWASP LLM, MITRE ATLAS — l'IA"),
        ("Ce qui nous distingue", "Trois choses",
         "1. Nos rapports sont mappés ATT&CK : votre SOC peut les exploiter\n"
         "   directement, au lieu de recevoir une liste de failles.\n\n"
         "2. Nous mesurons. ASVS, SAMM, CIS : des scores qui se re-mesurent.\n\n"
         "3. AI RedTeaming — une compétence rare, structurée sur OWASP LLM\n"
         "   et MITRE ATLAS."),
        ("Notre cadre", "Ce que nous exigeons avant de commencer",
         "NDA signé\nContrat et périmètre écrits\nRègles d'engagement validées\n"
         "Autorisation de test signée\n\n"
         "Ces exigences vous protègent autant qu'elles nous protègent."),
        ("Nos engagements", "Ce sur quoi vous pouvez compter",
         "Notification des vulnérabilités critiques sous 2 heures\n"
         "Preuves chiffrées, détruites à J+90, certificat à l'appui\n"
         "Aucune donnée sur un service tiers non contractualisé\n"
         "Un rapport relu par deux personnes avant de vous parvenir"),
        ("Parlons de vous", "", ""),
    ]
    for sur, titre, corps in plan:
        d = _diapo_vide(prs)
        _titre_diapo(d, prs, titre or sur, sur)
        if corps:
            _zone(d, 1.6, 3.0, 30, 14, corps, 14, False, C.TEXTE,
                  C.POLICE_CORPS, 1.35)
    prs.save(chemin)
    print("écrit", chemin)


def main() -> int:
    ap = argparse.ArgumentParser()
    ap.add_argument("--out", default="90-templates/build")
    a = ap.parse_args()
    os.makedirs(a.out, exist_ok=True)
    modele_rapport(os.path.join(a.out, "Modele-rapport.docx"))
    modele_restitution(os.path.join(a.out, "Modele-restitution.pptx"))
    modele_presentation(os.path.join(a.out, "Modele-presentation.pptx"))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
